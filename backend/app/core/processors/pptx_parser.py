"""PPTX 解析器：python-pptx 读幻灯片 → 形状文本 + 表格（表格渲染为 GFM Markdown）。"""
import io

from .base import BaseParser, register_parser, grid_to_markdown


@register_parser(".pptx")
class PptxParser(BaseParser):
    supported_ext = ".pptx"

    def extract(self, file_bytes: bytes) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            table_idx = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
                if shape.has_table:
                    table_idx += 1
                    rows = []
                    for row in shape.table.rows:
                        cells = []
                        prev = None
                        for cell in row.cells:
                            # python-pptx 合并单元格会重复返回同一对象，按对象去重
                            if cell is prev:
                                continue
                            prev = cell
                            cells.append(cell.text)
                        rows.append(cells)
                    md = grid_to_markdown(rows, f"### 表格 {table_idx}")
                    if md:
                        slide_texts.append(md)
            if slide_texts:
                parts.append(f"【第{slide_num}页】\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)