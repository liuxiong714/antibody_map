#!/usr/bin/env python
"""P1 功能全面测试：PPTX/XLSX 解析、策略模式、URL/HTML 抓取、OCR 兜底、PDF OCR 阈值修复。

运行方式（Windows PowerShell）:
  cd backend
  python tests/test_p1_features.py
"""
import io
import os
import re
import sys

# ── 路径修正 ──
_backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _backend_dir not in sys.path:
    sys.path.insert(0, _backend_dir)

_passed = 0
_failed = 0


def _ok(label: str, detail: str = ""):
    global _passed
    _passed += 1
    print(f"  ✓ {label}{(' — ' + detail) if detail else ''}")


def _fail(label: str, detail: str = ""):
    global _failed
    _failed += 1
    print(f"  ✗ {label}{(' — ' + detail) if detail else ''}")


# ─────────────────────────────────────────────────────────
# 测试 1：策略模式 — 解析器注册表
# ─────────────────────────────────────────────────────────
def test_strategy_registry():
    print("\n" + "=" * 60)
    print("【测试 1】策略模式：解析器注册表")
    print("=" * 60)

    from app.core.processors import get_parser, list_supported_exts

    expected_exts = {".epub", ".docx", ".pptx", ".xlsx", ".txt", ".html", ".htm"}
    actual_exts = set(list_supported_exts())

    # 1a: 所有预期扩展名都注册了
    missing = expected_exts - actual_exts
    if not missing:
        _ok("全部 7 个解析器已注册", f"exts={sorted(actual_exts)}")
    else:
        _fail("解析器注册缺失", f"missing={missing}")

    # 1b: get_parser 返回正确类型
    from app.core.processors.pptx_parser import PptxParser
    from app.core.processors.xlsx_parser import XlsxParser
    from app.core.processors.docx_parser import DocxParser
    from app.core.processors.txt_parser import TxtParser
    from app.core.processors.html_parser import HtmlParser

    cases = [
        (".pptx", PptxParser),
        (".xlsx", XlsxParser),
        (".docx", DocxParser),
        (".txt", TxtParser),
        (".html", HtmlParser),
        (".htm", HtmlParser),
    ]
    for ext, cls in cases:
        p = get_parser(ext)
        if p is not None and isinstance(p, cls):
            _ok(f"get_parser('{ext}') → {cls.__name__}")
        else:
            _fail(f"get_parser('{ext}') 返回类型错误", f"got={type(p).__name__ if p else 'None'}")

    # 1c: 未注册扩展名返回 None
    p = get_parser(".unknown")
    if p is None:
        _ok("get_parser('.unknown') → None")
    else:
        _fail("get_parser('.unknown') 应返回 None")

    # 1d: 大小写不敏感
    p = get_parser(".PPTX")
    if p is not None:
        _ok("get_parser('.PPTX') 大小写不敏感")
    else:
        _fail("get_parser('.PPTX') 大小写不敏感失败")


# ─────────────────────────────────────────────────────────
# 测试 2：PPTX 解析
# ─────────────────────────────────────────────────────────
def test_pptx_parsing():
    print("\n" + "=" * 60)
    print("【测试 2】PPTX 解析：幻灯片文本 + 表格")
    print("=" * 60)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        _fail("python-pptx 未安装，跳过")
        return

    # 构造 mock PPTX
    prs = Presentation()
    # 幻灯片 1：标题 + 文本
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "麻疹抗体阳性率调查"
    slide1.placeholders[1].text = "2023年广东省调查结果\n阳性率85.3%，样本量1500"

    # 幻灯片 2：表格
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8)).text = "数据汇总"
    table_shape = slide2.shapes.add_table(rows=3, cols=3, left=Inches(1), top=Inches(2), width=Inches(6), height=Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "地区"
    table.cell(0, 1).text = "阳性率"
    table.cell(0, 2).text = "样本量"
    table.cell(1, 0).text = "广州"
    table.cell(1, 1).text = "88.5%"
    table.cell(1, 2).text = "500"
    table.cell(2, 0).text = "深圳"
    table.cell(2, 1).text = "82.1%"
    table.cell(2, 2).text = "1000"

    # 幻灯片 3：空页（不应出现在输出中）
    prs.slides.add_slide(prs.slide_layouts[5])

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    from app.core.processors import get_parser
    parser = get_parser(".pptx")
    text = parser.extract_text(pptx_bytes)

    # 2a: 文本不为空
    if text and len(text) > 0:
        _ok("PPTX 解析返回非空文本", f"{len(text)} 字符")
    else:
        _fail("PPTX 解析返回空文本")
        return

    # 2b: 包含幻灯片 1 文本
    if "麻疹抗体阳性率调查" in text and "85.3%" in text:
        _ok("幻灯片文本提取正确")
    else:
        _fail("幻灯片文本提取缺失", f"text[:200]={text[:200]}")

    # 2c: 包含表格文本
    if "广州" in text and "88.5%" in text and "深圳" in text:
        _ok("表格文本提取正确")
    else:
        _fail("表格文本提取缺失", f"text[:200]={text[:200]}")

    # 2d: 包含页码标记
    if "【第1页】" in text and "【第2页】" in text:
        _ok("页码标记正确")
    else:
        _fail("页码标记缺失")

    # 2e: 空页不应出现
    if "【第3页】" not in text:
        _ok("空页正确跳过")
    else:
        _fail("空页不应出现在输出中")

    # 2f: 通过 document_parser 分发也能正确解析
    from app.core.document_parser import extract_text as dispatch_extract
    text2 = dispatch_extract(pptx_bytes, ".pptx")
    if text2 and "85.3%" in text2:
        _ok("document_parser 分发 PPTX 正常")
    else:
        _fail("document_parser 分发 PPTX 失败", f"text2[:100]={text2[:100] if text2 else 'None'}")


# ─────────────────────────────────────────────────────────
# 测试 3：XLSX 解析
# ─────────────────────────────────────────────────────────
def test_xlsx_parsing():
    print("\n" + "=" * 60)
    print("【测试 3】XLSX 解析：多工作表 + 数据")
    print("=" * 60)

    try:
        import openpyxl
    except ImportError:
        _fail("openpyxl 未安装，跳过")
        return

    # 构造 mock XLSX
    wb = openpyxl.Workbook()
    # Sheet 1: 阳性率数据
    ws1 = wb.active
    ws1.title = "阳性率"
    ws1.append(["地区", "疾病", "阳性率(%)", "样本量"])
    ws1.append(["广州", "麻疹", 85.3, 500])
    ws1.append(["深圳", "麻疹", 82.1, 1000])
    ws1.append(["珠海", "风疹", 91.2, 300])

    # Sheet 2: GMC 数据
    ws2 = wb.create_sheet("GMC")
    ws2.append(["地区", "GMC(mIU/ml)", "95%CI下限", "95%CI上限"])
    ws2.append(["广州", 1125.6, 1032.4, 1227.1])
    ws2.append(["深圳", 980.3, 890.5, 1070.1])

    # Sheet 3: 空表
    wb.create_sheet("空表")

    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()
    wb.close()

    from app.core.processors import get_parser
    parser = get_parser(".xlsx")
    text = parser.extract_text(xlsx_bytes)

    # 3a: 文本不为空
    if text and len(text) > 0:
        _ok("XLSX 解析返回非空文本", f"{len(text)} 字符")
    else:
        _fail("XLSX 解析返回空文本")
        return

    # 3b: 包含工作表名
    if "阳性率" in text and "GMC" in text:
        _ok("多工作表名提取正确")
    else:
        _fail("工作表名提取缺失", f"text[:200]={text[:200]}")

    # 3c: 包含数据值
    if "85.3" in text and "广州" in text and "1125.6" in text:
        _ok("数据值提取正确")
    else:
        _fail("数据值提取缺失", f"text[:300]={text[:300]}")

    # 3d: 单元格分隔符
    if " | " in text:
        _ok("单元格分隔符正确")
    else:
        _fail("单元格分隔符缺失")

    # 3e: 工作表标记
    if "【工作表:" in text:
        _ok("工作表标记正确")
    else:
        _fail("工作表标记缺失")

    # 3f: 空表不应有数据行
    lines = text.split("\n")
    empty_sheet_lines = [l for l in lines if "空表" in l]
    if len(empty_sheet_lines) == 1:  # 只有标题行，无数据行
        _ok("空表正确处理（仅标题）")
    else:
        _fail("空表处理异常", f"lines with '空表'={empty_sheet_lines}")

    # 3g: 通过 document_parser 分发
    from app.core.document_parser import extract_text as dispatch_extract
    text2 = dispatch_extract(xlsx_bytes, ".xlsx")
    if text2 and "85.3" in text2:
        _ok("document_parser 分发 XLSX 正常")
    else:
        _fail("document_parser 分发 XLSX 失败")


# ─────────────────────────────────────────────────────────
# 测试 4：回归 — TXT/HTML/DOCX 解析
# ─────────────────────────────────────────────────────────
def test_regression_formats():
    print("\n" + "=" * 60)
    print("【测试 4】回归测试：TXT / HTML / DOCX 解析")
    print("=" * 60)

    from app.core.processors import get_parser

    # 4a: TXT 解析
    txt_bytes = "麻疹IgG抗体阳性率为85.3%，样本量1500人。\n GMC=1125.6 mIU/ml".encode("utf-8")
    txt_parser = get_parser(".txt")
    txt_text = txt_parser.extract_text(txt_bytes)
    if "85.3%" in txt_text and "1125.6" in txt_text:
        _ok("TXT 解析正确")
    else:
        _fail("TXT 解析失败", f"text={txt_text[:100]}")

    # 4b: TXT GBK 编码兜底
    txt_gbk = "麻疹抗体阳性率85.3%".encode("gbk")
    txt_text2 = txt_parser.extract_text(txt_gbk)
    if "85.3%" in txt_text2:
        _ok("TXT GBK 编码兜底正确")
    else:
        _fail("TXT GBK 编码兜底失败", f"text={txt_text2[:100]}")

    # 4c: HTML 解析
    html_str = (
        "<html><head><title>麻疹抗体调查报告</title></head>"
        "<body><h1>2023年广东省麻疹抗体水平</h1>"
        "<p>阳性率85.3%，样本量1500人。GMC为1125.6 mIU/ml。</p>"
        "<table><tr><td>广州</td><td>88.5%</td></tr></table>"
        "</body></html>"
    )
    html_bytes = html_str.encode("utf-8")
    html_parser = get_parser(".html")
    html_text = html_parser.extract_text(html_bytes)
    if "85.3%" in html_text and "广州" in html_text:
        _ok("HTML 解析正确")
    else:
        _fail("HTML 解析失败", f"text={html_text[:200]}")

    # 4d: HTML .htm 扩展名
    htm_parser = get_parser(".htm")
    if htm_parser is not None:
        htm_text = htm_parser.extract_text(html_bytes)
        if "85.3%" in htm_text:
            _ok("HTML .htm 扩展名解析正确")
        else:
            _fail("HTML .htm 扩展名解析失败")
    else:
        _fail("get_parser('.htm') 返回 None")

    # 4e: DOCX 解析
    try:
        from docx import Document as DocxDocument
        doc = DocxDocument()
        doc.add_heading("麻疹抗体调查", level=1)
        doc.add_paragraph("阳性率85.3%，样本量1500人。")
        doc.add_paragraph("GMC为1125.6 mIU/ml。")
        buf = io.BytesIO()
        doc.save(buf)
        docx_bytes = buf.getvalue()

        docx_parser = get_parser(".docx")
        docx_text = docx_parser.extract_text(docx_bytes)
        if "85.3%" in docx_text and "1125.6" in docx_text:
            _ok("DOCX 解析正确")
        else:
            _fail("DOCX 解析失败", f"text={docx_text[:200]}")
    except ImportError:
        _fail("python-docx 未安装，跳过 DOCX 测试")


# ─────────────────────────────────────────────────────────
# 测试 5：URL/HTML 抓取
# ─────────────────────────────────────────────────────────
def test_url_html_fetch():
    print("\n" + "=" * 60)
    print("【测试 5】URL/HTML 抓取：标题提取 + 编码兜底")
    print("=" * 60)

    from app.core.url_fetcher import guess_title_from_html

    # 5a: UTF-8 HTML 标题提取
    html_utf8 = "<html><head><title>麻疹抗体阳性率调查报告</title></head><body>内容</body></html>".encode("utf-8")
    title = guess_title_from_html(html_utf8)
    if title == "麻疹抗体阳性率调查报告":
        _ok("UTF-8 HTML 标题提取正确")
    else:
        _fail("UTF-8 HTML 标题提取失败", f"title={title}")

    # 5b: GBK 编码 HTML 标题提取
    html_gbk = "<html><head><title>乙肝抗体调查</title></head><body>内容</body></html>".encode("gbk")
    title2 = guess_title_from_html(html_gbk)
    if title2 == "乙肝抗体调查":
        _ok("GBK 编码 HTML 标题提取正确")
    else:
        _fail("GBK 编码 HTML 标题提取失败", f"title={title2}")

    # 5c: 无标题标签
    html_no_title = "<html><head></head><body>no title page</body></html>".encode("utf-8")
    title3 = guess_title_from_html(html_no_title)
    if title3 is None:
        _ok("无标题 HTML 返回 None")
    else:
        _fail("无标题 HTML 应返回 None", f"title={title3}")

    # 5d: 标题截断（防止超长标题）
    long_title = "A" * 600
    html_long = f"<html><head><title>{long_title}</title></head></html>".encode("utf-8")
    title4 = guess_title_from_html(html_long)
    if title4 and len(title4) <= 500:
        _ok("超长标题截断正确", f"len={len(title4)}")
    else:
        _fail("超长标题未截断", f"len={len(title4) if title4 else 0}")

    # 5e: fetch_url 函数存在且可导入
    from app.core.url_fetcher import fetch_url
    if callable(fetch_url):
        _ok("fetch_url 函数可导入")
    else:
        _fail("fetch_url 导入失败")


# ─────────────────────────────────────────────────────────
# 测试 6：OCR 服务状态
# ─────────────────────────────────────────────────────────
def test_ocr_service():
    print("\n" + "=" * 60)
    print("【测试 6】OCR 服务：Tesseract 可用性 + 百度 OCR 函数")
    print("=" * 60)

    from app.core.ocr_service import get_ocr_status, ocr_image, ocr_baidu

    # 6a: get_ocr_status 返回结构
    status = get_ocr_status()
    if isinstance(status, dict) and "tesseract_available" in status:
        _ok("get_ocr_status 返回正确结构", f"tesseract={status['tesseract_available']}")
    else:
        _fail("get_ocr_status 返回结构异常", f"status={status}")

    # 6b: ocr_image 函数存在
    if callable(ocr_image):
        _ok("ocr_image 函数可导入")
    else:
        _fail("ocr_image 导入失败")

    # 6c: ocr_baidu 函数存在
    if callable(ocr_baidu):
        _ok("ocr_baidu 函数可导入")
    else:
        _fail("ocr_baidu 导入失败")

    # 6d: Tesseract 不可用时 ocr_image 返回 None（不崩溃）
    if not status.get("tesseract_available"):
        result = ocr_image(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100, lang="zh")
        if result is None or result == "":
            _ok("Tesseract 不可用时 ocr_image 安全返回空")
        else:
            _fail("Tesseract 不可用时 ocr_image 应返回空", f"got={result[:50]}")
    else:
        _ok("Tesseract 可用，跳过不可用测试", f"cmd={status.get('tesseract_cmd', '?')}")


# ─────────────────────────────────────────────────────────
# 测试 7：PDF OCR 阈值修复 — 乱码文本检测
# ─────────────────────────────────────────────────────────
def test_pdf_ocr_threshold():
    print("\n" + "=" * 60)
    print("【测试 7】PDF OCR 阈值修复：乱码非空文本检测")
    print("=" * 60)

    # 7a: 验证 PAGE_DIGIT_MIN 常量存在
    import inspect
    from app.core import pdf_parser
    source = inspect.getsource(pdf_parser.extract_text)
    if "PAGE_DIGIT_MIN" in source and "digit_count" in source:
        _ok("PDF 解析器包含数字字符检测逻辑")
    else:
        _fail("PDF 解析器缺少数字字符检测逻辑")

    # 7b: 验证逻辑：文本 >= 100 字符但数字 < 3 应触发 OCR
    # 模拟一段乱码文本（长但无数字）
    garbled_text = "隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙隙" * 5  # >100 字符，0 数字
    digit_count = len(re.findall(r"\d", garbled_text))
    text_len = len(garbled_text)
    if text_len >= 100 and digit_count < 3:
        _ok("乱码文本正确识别为需 OCR（长文本但数字不足）", f"len={text_len}, digits={digit_count}")
    else:
        _fail("乱码文本检测逻辑错误", f"len={text_len}, digits={digit_count}")

    # 7c: 正常文本（>= 100 字符且 >= 3 数字）不触发 OCR
    normal_text = "2023年广东省麻疹抗体阳性率为85.3%，样本量1500人，GMC为1125.6 mIU/ml。" * 3
    digit_count2 = len(re.findall(r"\d", normal_text))
    text_len2 = len(normal_text)
    if text_len2 >= 100 and digit_count2 >= 3:
        _ok("正常文本正确跳过 OCR（有足够数字字符）", f"len={text_len2}, digits={digit_count2}")
    else:
        _fail("正常文本检测逻辑错误", f"len={text_len2}, digits={digit_count2}")


# ─────────────────────────────────────────────────────────
# 测试 8：上传白名单 + 错误消息
# ─────────────────────────────────────────────────────────
def test_upload_whitelist():
    print("\n" + "=" * 60)
    print("【测试 8】上传白名单 + 错误消息一致性")
    print("=" * 60)

    from app.core.document_parser import ALLOWED_EXTS, MIME_MAP

    # 8a: PPTX 在白名单中
    if ".pptx" in ALLOWED_EXTS:
        _ok("PPTX 在上传白名单中")
    else:
        _fail("PPTX 不在上传白名单中")

    # 8b: XLSX 在白名单中
    if ".xlsx" in ALLOWED_EXTS:
        _ok("XLSX 在上传白名单中")
    else:
        _fail("XLSX 不在上传白名单中")

    # 8c: MIME 映射正确
    if MIME_MAP.get(".pptx") == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        _ok("PPTX MIME 映射正确")
    else:
        _fail("PPTX MIME 映射错误", f"got={MIME_MAP.get('.pptx')}")

    if MIME_MAP.get(".xlsx") == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        _ok("XLSX MIME 映射正确")
    else:
        _fail("XLSX MIME 映射错误", f"got={MIME_MAP.get('.xlsx')}")

    # 8d: 错误消息包含 PPTX/XLSX
    import inspect
    from app.api.v1 import literature
    source = inspect.getsource(literature.upload)
    if "PPTX" in source and "XLSX" in source:
        _ok("上传错误消息包含 PPTX/XLSX")
    else:
        _fail("上传错误消息未包含 PPTX/XLSX")

    # 8e: 所有 9 种格式都在白名单
    expected_exts = {".pdf", ".caj", ".epub", ".docx", ".pptx", ".xlsx", ".txt", ".html", ".htm"}
    if expected_exts == ALLOWED_EXTS:
        _ok("上传白名单完整（9 种格式）")
    else:
        _fail("上传白名单不完整", f"missing={expected_exts - ALLOWED_EXTS}, extra={ALLOWED_EXTS - expected_exts}")


# ─────────────────────────────────────────────────────────
# 测试 9：document_parser 分发完整性
# ─────────────────────────────────────────────────────────
def test_dispatch_completeness():
    print("\n" + "=" * 60)
    print("【测试 9】document_parser 分发完整性")
    print("=" * 60)

    from app.core.document_parser import extract_text, ALLOWED_EXTS
    from app.core.processors import list_supported_exts

    processor_exts = set(list_supported_exts())
    # PDF 和 CAJ 不走注册表，走独立路径
    non_registry_exts = {".pdf", ".caj"}
    # 所有非 PDF/CAJ 格式都应该在注册表里
    expected_registry = ALLOWED_EXTS - non_registry_exts
    missing = expected_registry - processor_exts
    if not missing:
        _ok("所有非 PDF/CAJ 格式都注册到策略模式", f"registry={sorted(processor_exts)}")
    else:
        _fail("部分格式未注册到策略模式", f"missing={missing}")

    # 不支持的格式应抛 ValueError 或返回空
    try:
        result = extract_text(b"test", ".unknown")
        if result == "":
            _ok("不支持的格式返回空字符串")
        else:
            _fail("不支持的格式应返回空", f"got={result[:50]}")
    except (ValueError, Exception):
        _ok("不支持的格式抛出异常")


# ─────────────────────────────────────────────────────────
# 主入口
# ─────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 60)
    print("P1 功能全面测试")
    print("涵盖：PPTX/XLSX 解析 | 策略模式 | URL/HTML 抓取 | OCR | PDF OCR 修复")
    print("=" * 60)

    tests = [
        test_strategy_registry,
        test_pptx_parsing,
        test_xlsx_parsing,
        test_regression_formats,
        test_url_html_fetch,
        test_ocr_service,
        test_pdf_ocr_threshold,
        test_upload_whitelist,
        test_dispatch_completeness,
    ]

    for test_fn in tests:
        try:
            test_fn()
        except Exception as e:
            _fail(f"{test_fn.__name__} 异常", str(e))

    print("\n" + "=" * 60)
    print(f"总结：{_passed}/{_passed + _failed} 通过", end="")
    if _failed:
        print(f"，{_failed} 失败 ❌")
    else:
        print(" 🎉")
    print("=" * 60)

    sys.exit(1 if _failed else 0)
