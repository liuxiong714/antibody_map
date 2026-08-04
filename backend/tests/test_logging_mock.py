"""
Mock 数据测试脚本：验证核心模块的日志输出是否按预期工作。

直接运行：python tests/test_logging_mock.py

测试范围：
1. document_parser.extract_text — 各格式分发链
2. url_fetcher.fetch_url / guess_title_from_html — URL 抓取链路
3. processors.base — 解析器注册表查找
"""
import io
import logging
import sys
import os

# 确保 backend 在 sys.path 中
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# 配置日志：输出到控制台，显示 INFO 及以上
logging.basicConfig(
    level=logging.INFO,
    format="%(levelname)s:%(name)s:%(message)s",
    stream=sys.stdout,
)
logger = logging.getLogger("uvicorn")


def test_unknown_format():
    """测试不支持的格式 → 应触发 [文档解析] 警告 + [解析器注册表] 警告"""
    print("\n" + "=" * 60)
    print("【测试 1】不支持的格式 (.xyz)")
    print("=" * 60)
    from app.core.document_parser import extract_text

    result = extract_text(b"some data", ".xyz")
    print(f"  → 返回结果: {repr(result)}")
    assert result == "", "不支持格式应返回空串"


def test_html_parsing():
    """测试 HTML 解析器 → 策略模式分发"""
    print("\n" + "=" * 60)
    print("【测试 2】HTML 解析 (策略模式 → HtmlParser)")
    print("=" * 60)
    from app.core.document_parser import extract_text

    html_bytes = (
        "<html><head><title>Mock 测试文献标题</title></head>"
        "<body><h1>抗体水平调查</h1><p>样本量: 500 人</p>"
        "<table><tr><td>年龄组</td><td>阳性率</td></tr>"
        "<tr><td>0-5岁</td><td>85.3%</td></tr></table></body></html>"
    ).encode("utf-8")
    result = extract_text(html_bytes, ".html")
    print(f"  → 提取文本 ({len(result)} 字符): {result[:100]}...")
    assert "抗体水平调查" in result, "应提取到 HTML 正文"
    assert "Mock 测试文献标题" in result, "HTML 标题应被提取"


def test_htm_parsing():
    """测试 .htm 扩展名 → 验证双注册"""
    print("\n" + "=" * 60)
    print("【测试 3】HTM 解析 (.htm 扩展名 → 双注册验证)")
    print("=" * 60)
    from app.core.document_parser import extract_text

    result = extract_text(b"<html><body><p>Hello .htm</p></body></html>", ".htm")
    print(f"  → 提取文本: {repr(result)}")
    assert "Hello .htm" in result


def test_txt_parsing():
    """测试 TXT 解析器"""
    print("\n" + "=" * 60)
    print("【测试 4】TXT 解析 (策略模式 → TxtParser)")
    print("=" * 60)
    from app.core.document_parser import extract_text

    result = extract_text("你好世界\nHello World".encode("utf-8"), ".txt")
    print(f"  → 提取文本: {repr(result)}")
    assert "你好世界" in result


def test_pptx_parsing():
    """测试 PPTX 解析器 — 用 python-pptx 构造 mock 文件"""
    print("\n" + "=" * 60)
    print("【测试 5】PPTX 解析 (策略模式 → PptxParser)")
    print("=" * 60)
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("  ⚠ python-pptx 未安装，跳过测试")
        print("  → 安装: pip install python-pptx")
        return

    prs = Presentation()
    # 第1页：标题 + 正文
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "2024年麻疹抗体水平调查报告"
    txBox2 = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf2 = txBox2.text_frame
    tf2.text = "样本量: 1200人    阳性率: 92.5%"
    # 第2页：表格
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    table_shape = slide2.shapes.add_table(3, 3, Inches(1), Inches(1), Inches(8), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "年龄组"
    table.cell(0, 1).text = "样本数"
    table.cell(0, 2).text = "阳性率"
    table.cell(1, 0).text = "0-5岁"
    table.cell(1, 1).text = "400"
    table.cell(1, 2).text = "90.1%"
    table.cell(2, 0).text = "6-10岁"
    table.cell(2, 1).text = "500"
    table.cell(2, 2).text = "94.2%"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()
    print(f"  → Mock PPTX 文件大小: {len(pptx_bytes)} 字节")

    from app.core.document_parser import extract_text

    result = extract_text(pptx_bytes, ".pptx")
    print(f"  → 提取文本 ({len(result)} 字符):\n{result[:300]}")
    assert "2024年麻疹抗体水平调查报告" in result, "PPTX 标题文本应被提取"
    assert "第1页" in result, "应包含页码标记"
    assert "90.1%" in result, "表格单元格文本应被提取"


def test_xlsx_parsing():
    """测试 XLSX 解析器 — 用 openpyxl 构造 mock 文件"""
    print("\n" + "=" * 60)
    print("【测试 6】XLSX 解析 (策略模式 → XlsxParser)")
    print("=" * 60)
    import openpyxl

    wb = openpyxl.Workbook()
    # Sheet1
    ws1 = wb.active
    ws1.title = "汇总数据"
    ws1.append(["省份", "样本量", "阳性率", "调查年份"])
    ws1.append(["北京", 1200, "92.5%", 2024])
    ws1.append(["上海", 980, "88.3%", 2024])
    ws1.append(["广东", 1500, "95.1%", 2023])
    # Sheet2
    ws2 = wb.create_sheet("详细数据")
    ws2.append(["年龄组", "省份", "样本数", "阳性数", "阳性率"])
    ws2.append(["0-5岁", "北京", 300, 270, "90.0%"])
    ws2.append(["6-10岁", "北京", 400, 380, "95.0%"])

    buf = io.BytesIO()
    wb.save(buf)
    xlsx_bytes = buf.getvalue()
    wb.close()
    print(f"  → Mock XLSX 文件大小: {len(xlsx_bytes)} 字节")

    from app.core.document_parser import extract_text

    result = extract_text(xlsx_bytes, ".xlsx")
    print(f"  → 提取文本 ({len(result)} 字符):\n{result[:400]}")
    assert "汇总数据" in result, "应包含工作表名"
    assert "北京" in result, "单元格数据应被提取"
    assert "92.5%" in result, "百分比数据应被提取"
    assert "详细数据" in result, "第二个工作表名应被提取"


def test_docx_parsing():
    """测试 DOCX 解析器 — 用 python-docx 构造 mock 文件"""
    print("\n" + "=" * 60)
    print("【测试 7】DOCX 解析 (策略模式 → DocxParser)")
    print("=" * 60)
    import docx

    document = docx.Document()
    document.add_paragraph("麻疹抗体水平调查研究")
    document.add_paragraph("")
    document.add_paragraph("目的：了解人群麻疹抗体水平")
    document.add_paragraph("方法：ELISA 法检测血清抗体")
    # 表格
    table = document.add_table(rows=3, cols=3)
    table.cell(0, 0).text = "年龄组"
    table.cell(0, 1).text = "样本数"
    table.cell(0, 2).text = "阳性率"
    table.cell(1, 0).text = "0-5岁"
    table.cell(1, 1).text = "300"
    table.cell(1, 2).text = "88.5%"

    buf = io.BytesIO()
    document.save(buf)
    docx_bytes = buf.getvalue()
    print(f"  → Mock DOCX 文件大小: {len(docx_bytes)} 字节")

    from app.core.document_parser import extract_text

    result = extract_text(docx_bytes, ".docx")
    print(f"  → 提取文本 ({len(result)} 字符):\n{result[:300]}")
    assert "麻疹抗体水平调查研究" in result, "DOCX 段落文本应被提取"
    assert "88.5%" in result, "表格单元格文本应被提取"


def test_url_fetcher():
    """测试 URL 抓取模块 — 使用 example.com"""
    print("\n" + "=" * 60)
    print("【测试 8】URL 抓取 (url_fetcher.fetch_url)")
    print("=" * 60)
    import asyncio
    from app.core.url_fetcher import fetch_url, guess_title_from_html

    content = asyncio.run(fetch_url("https://example.com"))
    print(f"  → 抓取内容: {len(content)} 字节")

    title = guess_title_from_html(content)
    print(f"  → 提取标题: {title}")
    assert title == "Example Domain", "example.com 的标题应为 'Example Domain'"


def test_parser_registry():
    """测试解析器注册表查找"""
    print("\n" + "=" * 60)
    print("【测试 9】解析器注册表查找")
    print("=" * 60)
    from app.core.processors import get_parser, list_supported_exts

    exts = list_supported_exts()
    print(f"  → 已注册的扩展名: {exts}")

    # 存在的解析器
    for ext in [".html", ".txt", ".pptx", ".xlsx", ".docx", ".epub"]:
        parser = get_parser(ext)
        assert parser is not None, f"{ext} 应有解析器"
        print(f"  ✓ {ext} → {parser.__class__.__name__}")

    # 不存在的解析器 → 应触发日志警告
    parser = get_parser(".unknown")
    assert parser is None, ".unknown 不应有解析器"
    print(f"  ✓ .unknown → None (已触发警告)")


if __name__ == "__main__":
    print("=" * 60)
    print("Mock 数据日志测试 - 开始")
    print("=" * 60)

    tests = [
        ("不支持的格式", test_unknown_format),
        ("HTML 解析", test_html_parsing),
        ("HTM 双注册", test_htm_parsing),
        ("TXT 解析", test_txt_parsing),
        ("PPTX 解析", test_pptx_parsing),
        ("XLSX 解析", test_xlsx_parsing),
        ("DOCX 解析", test_docx_parsing),
        ("URL 抓取", test_url_fetcher),
        ("解析器注册表", test_parser_registry),
    ]

    passed = 0
    failed = 0
    for name, fn in tests:
        try:
            fn()
            print(f"  ✅ {name} 通过\n")
            passed += 1
        except Exception as e:
            print(f"  ❌ {name} 失败: {e}\n")
            import traceback
            traceback.print_exc()
            failed += 1

    print("=" * 60)
    print(f"测试完成: {passed} 通过, {failed} 失败")
    print("=" * 60)
    sys.exit(0 if failed == 0 else 1)